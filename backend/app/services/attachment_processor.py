"""Attachment text extraction service supporting PDF, images (LLM vision), and Office formats.

Supports optional Presidio Image Redactor for PII detection/redaction directly on images.
Uses the agent's configured LLM (Axet / OpenAI / Azure — whatever is wired up) for
image analysis via LangChain multimodal messages.
"""

import io
import base64
from typing import Tuple, Optional

import structlog

logger = structlog.get_logger()

# Lazy-initialized Presidio image engines (expensive to create)
_image_redactor_engine = None
_image_analyzer_engine = None


def _get_presidio_image_engines():
    """Lazy-init Presidio image engines with Spanish NLP support."""
    global _image_redactor_engine, _image_analyzer_engine

    if _image_redactor_engine is not None:
        return _image_redactor_engine, _image_analyzer_engine

    try:
        from presidio_image_redactor import ImageRedactorEngine, ImageAnalyzerEngine
        from presidio_analyzer import AnalyzerEngine, PatternRecognizer, Pattern
        from presidio_analyzer.nlp_engine import SpacyNlpEngine

        # Use Spanish spaCy model for NER in images
        nlp_engine = SpacyNlpEngine(
            models=[{"lang_code": "es", "model_name": "es_core_news_lg"}],
        )
        analyzer = AnalyzerEngine(
            nlp_engine=nlp_engine,
            supported_languages=["es"],
        )

        # Register Spanish DNI/NIF/NIE/CIF + IBAN + phone + date recognizers.
        # The default PresidioAnalyzer ships English ones; we need Spanish forms
        # for OCR text extracted from DNI / payroll / invoice images.
        spanish_id_recognizer = PatternRecognizer(
            supported_entity="ES_DNI",
            supported_language="es",
            patterns=[
                Pattern(
                    "DNI_NIF",
                    r"(?i)(?:(?:DNI|NIF|NIE|NI|CIF)[:\s.\-]*)?\b\d{1,2}[\s.\-]?\d{3}[\s.\-]?\d{3}[\s.\-]?[A-Za-z]\b",
                    0.85,
                ),
                Pattern(
                    "NIE",
                    r"(?i)(?:(?:NIE|NI)[:\s.\-]*)?\b[XYZxyz][\s.\-]?\d{1,2}[\s.\-]?\d{3}[\s.\-]?\d{3}[\s.\-]?[A-Za-z]\b",
                    0.85,
                ),
                Pattern(
                    "CIF",
                    r"(?i)(?:CIF[:\s.\-]*)?\b[A-Ha-h]\d{7}[0-9A-Ja-j]\b",
                    0.80,
                ),
            ],
        )
        analyzer.registry.add_recognizer(spanish_id_recognizer)

        date_recognizer = PatternRecognizer(
            supported_entity="DATE_TIME",
            supported_language="es",
            patterns=[
                # Separators: / - . space | O (OCR misreads 0)
                Pattern("date_slash", r"\b\d{1,2}\s*[/\-.\\|]\s*\d{1,2}\s*[/\-.\\|]\s*\d{2,4}\b", 0.60),
                # Spaces (1 to 3) between groups: "10 01 1950"
                Pattern("date_spaced", r"\b\d{1,2}\s{1,3}\d{1,2}\s{1,3}\d{4}\b", 0.60),
                # With leading label: "FECHA NACIMIENTO 10 01 1950"
                Pattern("date_labelled", r"(?i)(?:fecha|nacim|nacido|expedic|valido|valid)[a-z\s]{0,20}\d{1,2}\s*\S?\s*\d{1,2}\s*\S?\s*\d{2,4}", 0.50),
            ],
        )
        analyzer.registry.add_recognizer(date_recognizer)

        # Spanish DNI support number: 3 letters + 6-7 digits (AAA023112).
        # OCR often splits or garbles it — allow whitespace/punct between groups.
        support_recognizer = PatternRecognizer(
            supported_entity="ES_DNI_SUPPORT",
            supported_language="es",
            patterns=[
                Pattern("support_num", r"\b[A-Z]{3}\s?\d{5,7}\b", 0.75),
                Pattern("support_ocr", r"\b[A-Z]{3,4}\s?\d{2,4}\s?\d{2,4}\b", 0.55),
            ],
        )
        analyzer.registry.add_recognizer(support_recognizer)

        # Single all-caps words that look like Spanish first names or surnames
        # on an ID. spaCy NER often skips short ALL-CAPS tokens — we help it.
        uppercase_name_recognizer = PatternRecognizer(
            supported_entity="PERSON",
            supported_language="es",
            patterns=[
                # 4+ consecutive uppercase letters (including Ñ, accented) — typical on scanned IDs.
                Pattern("allcaps_word", r"\b[A-ZÑÁÉÍÓÚ]{4,}\b", 0.45),
            ],
            # Avoid false positives on common ALL-CAPS labels
            deny_list=[
                "ESPAÑA", "ESPANA", "ESPAÑOLA", "DOCUMENTO", "NACIONAL",
                "IDENTIDAD", "NOMBRE", "APELLIDO", "APELLIDOS", "FECHA",
                "NACIMIENTO", "VALIDO", "HASTA", "SEXO", "NACIONALIDAD",
                "SOPORTE", "NUM", "NUMERO", "FIRMA", "DNI", "NIF",
            ],
        )
        analyzer.registry.add_recognizer(uppercase_name_recognizer)

        iban_recognizer = PatternRecognizer(
            supported_entity="IBAN_CODE",
            supported_language="es",
            patterns=[
                Pattern("iban_es", r"\bES\d{2}[\s]?\d{4}[\s]?\d{4}[\s]?\d{2}[\s]?\d{10}\b", 0.90),
            ],
        )
        analyzer.registry.add_recognizer(iban_recognizer)

        phone_recognizer = PatternRecognizer(
            supported_entity="PHONE_NUMBER",
            supported_language="es",
            patterns=[
                Pattern(
                    "phone_es",
                    r"(?<!\d)(?:\+34|0034)?\s?[6-9]\d{2}[\s.\-]?\d{3}[\s.\-]?\d{3}(?!\d)",
                    0.80,
                ),
            ],
        )
        analyzer.registry.add_recognizer(phone_recognizer)

        _image_analyzer_engine = ImageAnalyzerEngine(analyzer_engine=analyzer)
        _image_redactor_engine = ImageRedactorEngine(
            image_analyzer_engine=_image_analyzer_engine,
        )
        logger.info("presidio_image_engines_initialized")
        return _image_redactor_engine, _image_analyzer_engine
    except (ImportError, OSError) as e:
        logger.warning("presidio_image_not_available", error=str(e))
        return None, None


class AttachmentProcessor:
    """Extracts text from attachment files by routing on file extension.

    When presidio-image-redactor is installed, images can also be redacted
    (PII regions blacked out) and analyzed (PII entities with bounding boxes).
    """

    def extract_text(self, content: bytes, filename: str) -> Tuple[str, str]:
        """Extract text from attachment content.

        Args:
            content: Raw file bytes.
            filename: Original filename (used to determine format).

        Returns:
            Tuple of (extracted_text, format_used).
        """
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        if ext in ("jpg", "jpeg", "png", "bmp", "tiff", "tif"):
            return self._extract_image(content), "ocr"
        elif ext == "pdf":
            return self._extract_pdf(content), "pdf"
        elif ext == "docx":
            return self._extract_docx(content), "docx"
        elif ext == "xlsx":
            return self._extract_xlsx(content), "xlsx"
        elif ext == "pptx":
            return self._extract_pptx(content), "pptx"
        else:
            return self._extract_plaintext(content), "plaintext"

    # Tesseract struggles with small images: below this short-edge threshold,
    # we upscale + boost contrast + sharpen before running OCR.
    _MIN_OCR_SHORT_EDGE = 700   # below this → upscale up to 4x (cheaper than 6x Lanczos)
    _TARGET_OCR_SHORT_EDGE = 1600  # what we aim for after upscale

    def _prepare_ocr_image(self, content: bytes):
        """Open and preprocess an image for OCR quality.

        For tiny inputs we upscale with bicubic, convert to grayscale,
        apply autocontrast and a sharpen filter. On a 282x178 DNI JPG
        this takes OCR extraction from near-zero to extracting name,
        surnames, dates, support number and DNI number.

        Returns (pil_image, scale) where scale is the factor applied
        (1.0 means no upscale or post-processing).
        """
        from PIL import Image, ImageFilter, ImageOps
        image = Image.open(io.BytesIO(content))
        short = min(image.size)
        if short and short < self._MIN_OCR_SHORT_EDGE:
            scale = max(2.0, self._TARGET_OCR_SHORT_EDGE / short)
            new_size = (int(image.size[0] * scale), int(image.size[1] * scale))
            image = image.resize(new_size, Image.BICUBIC)
            # grayscale + autocontrast + sharpen → measured empirically
            # on a scanned DNI to be the best preprocessing for Tesseract.
            image = ImageOps.autocontrast(ImageOps.grayscale(image))
            image = image.filter(ImageFilter.SHARPEN)
            # Convert back to RGB so Presidio ImageRedactor can draw RGB
            # fill rectangles over it (grayscale mode breaks its color arg).
            image = image.convert("RGB")
            return image, scale
        # Ensure mode is compatible with redactor (some TIFFs come as "L" or "P")
        if image.mode != "RGB":
            image = image.convert("RGB")
        return image, 1.0

    # Tesseract page-segmentation mode: 6 = "assume a single uniform block
    # of text". Works best on ID/scanned documents. Auto (PSM 3) tends to
    # skip small fields like name/date on ID cards.
    _OCR_CONFIG = "--psm 6"

    def redact_image(self, content: bytes, fill_color: Tuple[int, int, int] = (0, 0, 0)) -> Optional[bytes]:
        """Redact PII from an image using Presidio Image Redactor.

        Args:
            content: Raw image bytes (JPEG, PNG, etc.)
            fill_color: RGB color to fill redacted regions (default: black)

        Returns:
            Redacted image as PNG bytes, or None if Presidio is not available.
        """
        redactor, _ = _get_presidio_image_engines()
        if redactor is None:
            return None

        try:
            image, scale = self._prepare_ocr_image(content)
            redacted = redactor.redact(
                image,
                fill=fill_color,
                ocr_kwargs={"lang": "spa+eng", "config": self._OCR_CONFIG},
                language="es",
            )
            buf = io.BytesIO()
            redacted.save(buf, format="PNG")
            logger.info("image_redacted_successfully", scale=scale, size=image.size)
            return buf.getvalue()
        except Exception as e:
            logger.error("image_redaction_failed", error=str(e))
            return None

    def analyze_image(self, content: bytes) -> list:
        """Analyze an image for PII entities using Presidio Image Analyzer.

        Args:
            content: Raw image bytes.

        Returns:
            List of detected PII entities with bounding boxes,
            or empty list if Presidio is not available.
        """
        _, analyzer = _get_presidio_image_engines()
        if analyzer is None:
            return []

        try:
            image, scale = self._prepare_ocr_image(content)
            results = analyzer.analyze(
                image,
                ocr_kwargs={"lang": "spa+eng", "config": self._OCR_CONFIG},
                language="es",
            )
            entities = []
            for r in results:
                entities.append({
                    "entity_type": r.entity_type,
                    "text": r.text if hasattr(r, "text") else "",
                    "score": r.score,
                    "left": r.left,
                    "top": r.top,
                    "width": r.width,
                    "height": r.height,
                })
            logger.info("image_analyzed", entities_found=len(entities), scale=scale)
            return entities
        except Exception as e:
            logger.error("image_analysis_failed", error=str(e))
            return []

    # Prompt que instruye al LLM a analizar la imagen SIN devolver PII
    IMAGE_ANALYSIS_PROMPT = """Analiza esta imagen y describe su contenido.

REGLAS CRITICAS DE PRIVACIDAD — NUNCA incluyas en tu respuesta:
- Nombres de personas (reemplazar por [NOMBRE])
- Numeros de documento, DNI, pasaporte, NIE (reemplazar por [DOCUMENTO])
- Fechas de nacimiento (reemplazar por [FECHA_NACIMIENTO])
- Direcciones postales (reemplazar por [DIRECCION])
- Numeros de telefono (reemplazar por [TELEFONO])
- Emails (reemplazar por [EMAIL])
- Numeros de cuenta, tarjeta de credito (reemplazar por [CUENTA])
- Matriculas de vehiculos (reemplazar por [MATRICULA])
- Cualquier otro dato que identifique a una persona (reemplazar por [PII])

Tu respuesta debe incluir:
1. TIPO DE DOCUMENTO: Que tipo de imagen/documento es (DNI, pasaporte, factura, captura de pantalla, foto, etc.)
2. PII DETECTADO: Lista de tipos de datos personales encontrados (sin revelar los valores)
3. CONTENIDO SEGURO: Descripcion del contenido con todos los datos personales reemplazados por los tokens indicados arriba

NUNCA reveles el valor real de ningun dato personal. Esta es una aplicacion GDPR-compliant."""

    def _extract_image(self, content: bytes) -> str:
        """Sync wrapper kept for backwards compatibility (e.g. _extract_pdf
        fallback). Runs the async analyzer in a fresh loop when no loop is
        running; otherwise delegates to the caller via extract_text_async.
        """
        import asyncio
        try:
            asyncio.get_running_loop()
            # We are already inside an event loop — caller should use
            # extract_text_async (awaitable) instead. Block-invoke would
            # deadlock, so return a safe placeholder.
            return "[Analisis de imagen requiere contexto async — usa extract_text_async]"
        except RuntimeError:
            return asyncio.run(self._extract_image_async(content))

    async def _extract_image_async(self, content: bytes) -> str:
        """Analyze image using the agent's configured LLM (Axet) asynchronously.

        Uses `await llm.ainvoke(...)` so FastAPI's event loop is not blocked
        during the (potentially multi-second) LLM Vision call.
        """
        try:
            from ..main import app_state
            from langchain_core.messages import HumanMessage

            agent = app_state.get("agent")
            if agent is None or getattr(agent, "llm", None) is None:
                return "[Analisis de imagenes no disponible: agente LLM no inicializado]"

            b64_image = base64.b64encode(content).decode("utf-8")

            mime = "image/png"
            if content[:3] == b'\xff\xd8\xff':
                mime = "image/jpeg"
            elif content[:4] == b'\x89PNG':
                mime = "image/png"
            elif content[:2] == b'BM':
                mime = "image/bmp"

            message = HumanMessage(content=[
                {"type": "text", "text": self.IMAGE_ANALYSIS_PROMPT},
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64_image}"}},
            ])

            response = await agent.llm.ainvoke([message])
            text = (response.content or "").strip()
            logger.info(
                "image_analyzed_via_agent_llm",
                length=len(text),
                provider=getattr(agent, "_llm_provider", "unknown"),
                model=getattr(agent, "_llm_model", "unknown"),
            )
            return text

        except Exception as e:
            logger.error("llm_image_analysis_failed", error=str(e))
            return f"[Error analisis imagen: {e}]"

    async def extract_text_async(self, content: bytes, filename: str) -> Tuple[str, str]:
        """Async version of extract_text — use from async contexts to avoid
        blocking the event loop while LLM Vision processes an image.
        """
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext in ("jpg", "jpeg", "png", "bmp", "tiff", "tif"):
            text = await self._extract_image_async(content)
            return text, "ocr"
        # Other formats are CPU-bound but short; run in a thread to keep the
        # event loop responsive on large PDFs / xlsx files.
        import asyncio
        text, fmt = await asyncio.to_thread(self.extract_text, content, filename)
        return text, fmt

    def _extract_pdf(self, content: bytes) -> str:
        try:
            import pdfplumber
        except ImportError:
            return "[Error: pdfplumber no instalado. Instalar con: pip install pdfplumber]"

        try:
            pages_text = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        pages_text.append(text.strip())
                    else:
                        # LLM vision fallback for scanned pages
                        img = page.to_image(resolution=300)
                        buf = io.BytesIO()
                        img.original.save(buf, format="PNG")
                        ocr_text = self._extract_image(buf.getvalue())
                        if ocr_text and not ocr_text.startswith("[Error"):
                            pages_text.append(ocr_text)

            return "\n\n".join(pages_text) if pages_text else "[PDF sin texto extraíble]"
        except Exception as e:
            logger.error("pdf_extraction_failed", error=str(e))
            return f"[Error PDF: {e}]"

    def _extract_docx(self, content: bytes) -> str:
        try:
            from docx import Document
        except ImportError:
            return "[Error: python-docx no instalado. Instalar con: pip install python-docx]"

        try:
            doc = Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs) if paragraphs else "[DOCX sin texto]"
        except Exception as e:
            logger.error("docx_extraction_failed", error=str(e))
            return f"[Error DOCX: {e}]"

    def _extract_xlsx(self, content: bytes) -> str:
        try:
            from openpyxl import load_workbook
        except ImportError:
            return "[Error: openpyxl no instalado. Instalar con: pip install openpyxl]"

        try:
            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            rows_text = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows_text.append(f"--- Hoja: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows_text.append(" | ".join(cells))
            wb.close()
            return "\n".join(rows_text) if rows_text else "[XLSX sin datos]"
        except Exception as e:
            logger.error("xlsx_extraction_failed", error=str(e))
            return f"[Error XLSX: {e}]"

    def _extract_pptx(self, content: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return "[Error: python-pptx no instalado. Instalar con: pip install python-pptx]"

        try:
            prs = Presentation(io.BytesIO(content))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                texts.append(para.text.strip())
                if texts:
                    slides_text.append(f"--- Diapositiva {i} ---\n" + "\n".join(texts))
            return "\n\n".join(slides_text) if slides_text else "[PPTX sin texto]"
        except Exception as e:
            logger.error("pptx_extraction_failed", error=str(e))
            return f"[Error PPTX: {e}]"

    def _extract_plaintext(self, content: bytes) -> str:
        try:
            return content.decode("utf-8").strip()
        except UnicodeDecodeError:
            try:
                return content.decode("latin-1").strip()
            except Exception:
                return "[No se pudo decodificar el archivo como texto]"
